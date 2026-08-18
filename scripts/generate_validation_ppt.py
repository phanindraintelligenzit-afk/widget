import collections
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

def create_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # Color Palette - Sleek Corporate Dark Theme
    c_dark = RGBColor(9, 13, 22)         # Deep Navy Background
    c_slate = RGBColor(30, 41, 59)       # Card background
    c_green = RGBColor(16, 185, 129)     # Emerald Green Accent
    c_blue = RGBColor(56, 189, 248)      # Light Sky Blue Accent
    c_white = RGBColor(255, 255, 255)    # White text
    c_gray = RGBColor(148, 163, 184)     # Muted text
    c_red = RGBColor(239, 68, 68)        # Error/Alert Red

    artifact_dir = r"C:\Users\User\.gemini\antigravity\brain\d42c9526-0eea-40b3-97d6-ee5b2265af2c"

    def add_solid_slide():
        s = prs.slides.add_slide(blank_layout)
        bg = s.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = c_dark
        return s

    def add_slide_header(slide, title_text, category_text="DPI-LS PERFORMANCE INDEX"):
        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.4))
        p_cat = cat_box.text_frame.paragraphs[0]
        p_cat.text = category_text.upper()
        p_cat.font.name = "Arial"
        p_cat.font.size = Pt(10)
        p_cat.font.bold = True
        p_cat.font.color.rgb = c_green
        
        t_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11.7), Inches(0.8))
        p = t_box.text_frame.paragraphs[0]
        p.text = title_text
        p.font.name = "Arial"
        p.font.size = Pt(26)
        p.font.bold = True
        p.font.color.rgb = c_white

    # -------------------------------------------------------------
    # SLIDE 1: Title
    # -------------------------------------------------------------
    s1 = add_solid_slide()
    accent_bar = s1.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.4), Inches(7.5))
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = c_green
    accent_bar.line.fill.background()

    t_box = s1.shapes.add_textbox(Inches(1.2), Inches(2.0), Inches(11.0), Inches(2.2))
    tf = t_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "DIGITAL FTE"
    p.font.name = "Arial"
    p.font.size = Pt(56)
    p.font.bold = True
    p.font.color.rgb = c_white
    
    p2 = tf.add_paragraph()
    p2.text = "PERFORMANCE INDEX – VALIDATION DIMENSION"
    p2.font.name = "Arial"
    p2.font.size = Pt(28)
    p2.font.bold = True
    p2.font.color.rgb = c_blue
    p2.space_before = Pt(10)

    sub_box = s1.shapes.add_textbox(Inches(1.2), Inches(4.5), Inches(11.0), Inches(2.0))
    tf_sub = sub_box.text_frame
    tf_sub.word_wrap = True
    p3 = tf_sub.paragraphs[0]
    p3.text = "Cost Traceability, Real-Time Telemetry & Core Scoring Validation"
    p3.font.name = "Arial"
    p3.font.size = Pt(16)
    p3.font.color.rgb = c_gray
    
    p4 = tf_sub.add_paragraph()
    p4.text = "Presented to: Wipro, Deloitte, JP Morgan, EY, PwC, Infosys & Regeneron"
    p4.font.name = "Arial"
    p4.font.size = Pt(14)
    p4.font.bold = True
    p4.font.color.rgb = c_green
    p4.space_before = Pt(8)

    p5 = tf_sub.add_paragraph()
    p5.text = "Executive Briefing: CEO, Ranga (CEO), Phani (Lead)"
    p5.font.name = "Arial"
    p5.font.size = Pt(12)
    p5.font.color.rgb = c_white
    p5.space_before = Pt(5)

    p6 = tf_sub.add_paragraph()
    p6.text = "Status: Development Complete\nProduction Validation Pending (Regeneron AWS/RDS + Approved Benchmark Dataset)"
    p6.font.name = "Arial"
    p6.font.size = Pt(12)
    p6.font.color.rgb = c_red
    p6.font.bold = True
    p6.space_before = Pt(15)

    # -------------------------------------------------------------
    # SLIDE 2: What is DPI-LS Validation
    # -------------------------------------------------------------
    s2 = add_solid_slide()
    add_slide_header(s2, "What is DPI-LS Validation & How It Works")
    
    tb_left = s2.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(5.5), Inches(5.0))
    tf_l = tb_left.text_frame
    tf_l.word_wrap = True
    
    p = tf_l.paragraphs[0]
    p.text = "Digital FTE Validation Index"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = c_blue
    p.space_after = Pt(12)
    
    bullets = [
        "Validation (V) measures whether an AI agent produces trustworthy, accurate, observable and verifiable outputs during runtime.",
        "The Validation dimension combines three runtime validation resources:\n  • DeepEval\n  • Jaeger\n  • Zipkin",
        "Together they evaluate:\n  • AI answer quality\n  • Runtime correctness\n  • Distributed execution tracing\n  • End-to-end observability",
        "Validation ensures every AI execution can be evaluated, traced and verified before production deployment.",
        "Runtime telemetry is collected from DeepEval, Jaeger, and Zipkin using real OpenTelemetry traces in the development environment."
    ]
    for b in bullets:
        p = tf_l.add_paragraph()
        p.text = "• " + b
        p.font.size = Pt(13)
        p.font.color.rgb = c_gray
        p.space_after = Pt(10)

    # Use uploaded screenshot for validation
    val_img_path = os.path.join(artifact_dir, "media__1782888245373.png") # DeepEval is a good observability illustration
    if os.path.exists(val_img_path):
        s2.shapes.add_picture(val_img_path, Inches(6.8), Inches(1.8), Inches(5.7), Inches(4.5))

    # -------------------------------------------------------------
    # SLIDE 3: Telemetry Sources
    # -------------------------------------------------------------
    s3 = add_solid_slide()
    add_slide_header(s3, "Validation (V) Dimension & Telemetry Sources")
    
    desc_box = s3.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.7), Inches(1.2))
    tf_desc = desc_box.text_frame
    tf_desc.word_wrap = True
    p = tf_desc.paragraphs[0]
    p.text = "Runtime Metrics"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = c_blue
    p.space_after = Pt(6)
    
    p2 = tf_desc.add_paragraph()
    p2.text = "DeepEval: Answer Relevancy, Faithfulness, Hallucination, Correctness | Jaeger: Trace ID, Span Count, Latency, Dependencies | Zipkin: Trace Timeline, Execution Timeline, Service Calls, Trace Latency"
    p2.font.size = Pt(12)
    p2.font.color.rgb = c_gray

    rows, cols = 4, 5
    left, top, width, height = Inches(0.8), Inches(2.9), Inches(11.7), Inches(3.8)
    table_shape = s3.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    
    table.columns[0].width = Inches(0.5)
    table.columns[1].width = Inches(1.8)
    table.columns[2].width = Inches(3.2)
    table.columns[3].width = Inches(2.2)
    table.columns[4].width = Inches(4.0)
    
    headers = ["No", "Parameter", "What It Measures", "Source System", "URL / Where Data Comes From"]
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = c_slate
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = c_white
        
    data = [
        ["1", "LLM Quality Evaluation", "Evaluates Answer Relevancy, Faithfulness, Hallucination and Correctness.", "DeepEval SDK", "https://deepeval.com"],
        ["2", "Distributed Trace Monitoring", "Captures runtime spans, latency, service dependencies and execution flow.", "Jaeger", "http://localhost:16686"],
        ["3", "Execution Timeline Visualization", "Captures request timeline, trace visualization, service execution path and runtime analysis.", "Zipkin", "http://localhost:9411"]
    ]
    for row_idx, row_val in enumerate(data):
        for col_idx, text in enumerate(row_val):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = text
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(15, 23, 42)
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(10)
            p.font.color.rgb = c_gray

    # -------------------------------------------------------------
    # SLIDE 4: Formula & Examples
    # -------------------------------------------------------------
    s4 = add_solid_slide()
    add_slide_header(s4, "Formula, Examples & Validation Rule")
    
    examples_img_path = os.path.join(artifact_dir, "media__1782888216173.png")
    if os.path.exists(examples_img_path):
        s4.shapes.add_picture(examples_img_path, Inches(0.8), Inches(1.5), Inches(7.5), Inches(5.2))
        
    expl_card = s4.shapes.add_textbox(Inches(8.5), Inches(1.5), Inches(4.0), Inches(5.2))
    tf_expl = expl_card.text_frame
    tf_expl.word_wrap = True
    
    p = tf_expl.paragraphs[0]
    p.text = "Key Logic Highlights"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = c_blue
    p.space_after = Pt(12)
    
    bullets_expl = [
        "Formula:\nValidation Score = Correct Validation Checks ÷ Total Validation Checks",
        "The validation score represents how many required validation checks passed successfully.",
        "If all validation checks pass, Validation Score = 1.0",
        "If several checks fail, Validation Score decreases accordingly.",
        "Good Example:\nPassed = 10 / Required = 10\nScore = 1.0 (Runtime successfully validated)",
        "Bad Example:\nPassed = 5 / Required = 10\nScore = 0.5 (Validation confidence decreases)",
        "Safety Rule: If Validation Score drops below the configured threshold, the overall DPI-LS score is reduced and the dashboard flags the execution for investigation."
    ]
    for b in bullets_expl:
        p = tf_expl.add_paragraph()
        p.text = "• " + b
        p.font.size = Pt(12)
        p.font.color.rgb = c_gray
        p.space_after = Pt(8)

    # -------------------------------------------------------------
    # SLIDE 5: Dashboards
    # -------------------------------------------------------------
    s5 = add_solid_slide()
    add_slide_header(s5, "DPI-LS Validation Dashboards & Runtime Evidence")
    
    # Left: Resources Dashboard (I'll just reuse resources_page_1782405844627.png as it shows the Resources Dashboard)
    res_img_path = os.path.join(artifact_dir, "resources_page_1782405844627.png")
    if os.path.exists(res_img_path):
        s5.shapes.add_picture(res_img_path, Inches(0.8), Inches(1.5), height=Inches(4.5))
        lbl_res = s5.shapes.add_textbox(Inches(0.8), Inches(6.0), Inches(5.0), Inches(0.4))
        lbl_res.text_frame.paragraphs[0].text = "Validation Resource Evaluation Dashboard"
        lbl_res.text_frame.paragraphs[0].font.size = Pt(11)
        lbl_res.text_frame.paragraphs[0].font.color.rgb = c_gray
        lbl_res.text_frame.paragraphs[0].font.italic = True

    # Right: 3 stacked screenshots
    de_img = os.path.join(artifact_dir, "media__1782888245373.png") # DeepEval
    j_img = os.path.join(artifact_dir, "media__1782888234232.png") # Jaeger
    z_img = os.path.join(artifact_dir, "media__1782888224572.png") # Zipkin
    
    y_offsets = [1.5, 3.0, 4.5]
    caps = ["LLM Evaluation Metrics", "Distributed Runtime Tracing", "Trace Timeline Visualization"]
    for i, (img, cap) in enumerate(zip([de_img, j_img, z_img], caps)):
        if os.path.exists(img):
            s5.shapes.add_picture(img, Inches(7.5), Inches(y_offsets[i]), width=Inches(5.0))
            lbl = s5.shapes.add_textbox(Inches(7.5), Inches(y_offsets[i] + 1.2), Inches(5.0), Inches(0.3))
            lbl.text_frame.paragraphs[0].text = cap
            lbl.text_frame.paragraphs[0].font.size = Pt(10)
            lbl.text_frame.paragraphs[0].font.color.rgb = c_gray
            lbl.text_frame.paragraphs[0].font.italic = True

    footer = s5.shapes.add_textbox(Inches(0.8), Inches(6.8), Inches(11.7), Inches(0.4))
    footer.text_frame.paragraphs[0].text = "Validation dashboards showing real runtime evaluation and distributed observability.\nEvidence shown is from live runtime telemetry captured in the development environment using OpenTelemetry."
    footer.text_frame.paragraphs[0].font.size = Pt(11)
    footer.text_frame.paragraphs[0].font.color.rgb = c_white

    # -------------------------------------------------------------
    # SLIDE 6: Workflow
    # -------------------------------------------------------------
    s6 = add_solid_slide()
    add_slide_header(s6, "Production Validation Execution Workflow")
    
    tb_cmd = s6.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.7), Inches(5.2))
    tf_c = tb_cmd.text_frame
    tf_c.word_wrap = True
    
    cmds = [
        ["Step 1: Start Jaeger", "Docker container (Port 16686)"],
        ["Step 2: Start Zipkin", "Docker container (Port 9411)"],
        ["Step 3: Start Backend", "uv run uvicorn api.app:app --host 127.0.0.1 --port 8000 --reload"],
        ["Step 4: Run Validation Agent", "uv run python examples/test_agent.py"],
        ["Step 5: Execute Validation", "DeepEval evaluation, Jaeger tracing, Zipkin tracing"],
        ["Step 6: Verify Results", "Validation Dashboard, DeepEval Metrics, Jaeger Trace, Zipkin Timeline, Resources Dashboard"]
    ]
    for title, cmd_str in cmds:
        p = tf_c.add_paragraph()
        p.text = "■ " + title
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = c_green
        p.space_after = Pt(2)
        
        p2 = tf_c.add_paragraph()
        p2.text = cmd_str
        p2.font.name = "Consolas"
        p2.font.size = Pt(11)
        p2.font.color.rgb = c_blue
        p2.space_after = Pt(8)

    p_exp = tf_c.add_paragraph()
    p_exp.text = "\nExpected Results:\nDeepEval → SUCCESS\nJaeger → SUCCESS\nZipkin → SUCCESS\nResources Dashboard → SUCCESS\nLive Runtime Metrics Displayed"
    p_exp.font.size = Pt(14)
    p_exp.font.bold = True
    p_exp.font.color.rgb = c_white
    p_exp.space_after = Pt(6)

    # -------------------------------------------------------------
    # SLIDE 7: Summary
    # -------------------------------------------------------------
    s7 = add_solid_slide()
    add_slide_header(s7, "Validation Dimension Summary")
    
    tb_sum = s7.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.7), Inches(5.2))
    tf_s = tb_sum.text_frame
    tf_s.word_wrap = True
    
    bullets_sum = [
        "✅ DeepEval integrated",
        "✅ Jaeger integrated",
        "✅ Zipkin integrated",
        "✅ OpenTelemetry implemented",
        "✅ Validation Resources Dashboard",
        "✅ Live Runtime Metrics",
        "✅ FastAPI Backend",
        "✅ SQLite Integration",
        "✅ Development Testing Completed"
    ]
    for b in bullets_sum:
        p = tf_s.add_paragraph()
        p.text = b
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = c_white
        p.space_after = Pt(6)

    p_next = tf_s.add_paragraph()
    p_next.text = "\nNext Steps:"
    p_next.font.size = Pt(18)
    p_next.font.bold = True
    p_next.font.color.rgb = c_blue
    p_next.space_after = Pt(8)

    next_steps = [
        "Regeneron AWS validation",
        "Regeneron RDS",
        "Approved Benchmark Dataset",
        "Production Recall Measurement",
        "M6 Validation"
    ]
    for n in next_steps:
        p = tf_s.add_paragraph()
        p.text = "• " + n
        p.font.size = Pt(14)
        p.font.color.rgb = c_gray
        p.space_after = Pt(6)

    prs.save("dpi-ls-validation.pptx")
    print("PowerPoint file dpi-ls-validation.pptx created successfully.")

if __name__ == "__main__":
    create_deck()
