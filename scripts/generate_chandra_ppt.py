from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    prs = Presentation()

    # Define common layout indices (standard PPT layouts)
    TITLE_SLIDE_LAYOUT = 0
    BULLET_SLIDE_LAYOUT = 1
    TITLE_ONLY_LAYOUT = 5

    # -----------------------------------------
    # Slide 1: Title Slide
    # -----------------------------------------
    slide1 = prs.slides.add_slide(prs.slide_layouts[TITLE_SLIDE_LAYOUT])
    title1 = slide1.shapes.title
    subtitle1 = slide1.placeholders[1]

    title1.text = "Chandra + DPI: Task In, Trusted Score Out"
    subtitle1.text = "Worker • Witness • Judge\nThree separate parties, by design."

    # Customizing fonts
    title1.text_frame.paragraphs[0].font.bold = True
    title1.text_frame.paragraphs[0].font.size = Pt(40)
    title1.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)

    # -----------------------------------------
    # Slide 2: Why 3rd Party Resources?
    # -----------------------------------------
    slide2 = prs.slides.add_slide(prs.slide_layouts[BULLET_SLIDE_LAYOUT])
    title2 = slide2.shapes.title
    body2 = slide2.placeholders[1]

    title2.text = "Why Do We Need Independent Resources?"
    
    tf2 = body2.text_frame
    tf2.text = "Chandra is the Worker. DPI-LS is the Appraisal System."
    
    p = tf2.add_paragraph()
    p.text = "The Problem of Self-Bias:"
    p.level = 1
    p = tf2.add_paragraph()
    p.text = "An AI cannot objectively grade its own homework. Asking an AI if it hallucinated yields unreliable results."
    p.level = 2
    
    p = tf2.add_paragraph()
    p.text = "AWS APIs are NOT 'Third Party':"
    p.level = 1
    p = tf2.add_paragraph()
    p.text = "They are our direct window into the client's cloud. Chandra knows nothing in advance; every fact is pulled fresh at task time."
    p.level = 2

    p = tf2.add_paragraph()
    p.text = "Why Independent Witnesses?"
    p.level = 1
    p = tf2.add_paragraph()
    p.text = "Without external observability (DeepEval, Zipkin, LangSmith), the DPI score means nothing to a buyer. Independence makes the score sellable."
    p.level = 2

    # -----------------------------------------
    # Slide 3: Chandra — The Worker Workflow
    # -----------------------------------------
    slide3 = prs.slides.add_slide(prs.slide_layouts[BULLET_SLIDE_LAYOUT])
    title3 = slide3.shapes.title
    body3 = slide3.placeholders[1]

    title3.text = "Chandra — The Worker (5-Step Fixed Procedure)"
    
    tf3 = body3.text_frame
    tf3.text = "The Brain: Claude on AWS Bedrock (data never leaves AWS) | The Toolbelt: AWS APIs (boto3)"
    tf3.paragraphs[0].font.bold = True
    
    p = tf3.add_paragraph()
    p.text = "1. Classify: What kind of problem is this? (Cost, Security, Compliance, etc.)"
    p.level = 1
    
    p = tf3.add_paragraph()
    p.text = "2. Investigate: Pull actual facts from AWS APIs. Real numbers, no guessing."
    p.level = 1

    p = tf3.add_paragraph()
    p.text = "3. Assess Blast Radius (The Judgment Step):"
    p.level = 1
    p = tf3.add_paragraph()
    p.text = "Before acting, what breaks? E.g., refusing to kill an expensive GxP-validated genomics workload to save money. This is judgment, not automation."
    p.level = 2

    p = tf3.add_paragraph()
    p.text = "4. Recommend: Propose action + full reasoning (the safe move, not just the cheap one)."
    p.level = 1

    p = tf3.add_paragraph()
    p.text = "5. Audit Summary: Writes everything it saw, thought, and decided into the work register (the Witness)."
    p.level = 1

    # -----------------------------------------
    # Slide 4: DPI — The Judge
    # -----------------------------------------
    slide4 = prs.slides.add_slide(prs.slide_layouts[BULLET_SLIDE_LAYOUT])
    title4 = slide4.shapes.title
    body4 = slide4.placeholders[1]

    title4.text = "DPI — The Judge (Appraisal System)"
    
    tf4 = body4.text_frame
    tf4.text = "DPI never asks Chandra how well it did. It reads independent evidence only."
    tf4.paragraphs[0].font.bold = True
    
    p = tf4.add_paragraph()
    p.text = "Feed 1: Work Register"
    p.level = 1
    p = tf4.add_paragraph()
    p.text = "Every task, decision, and outcome Chandra logged (Feeds P, Q, E, G, R)."
    p.level = 2

    p = tf4.add_paragraph()
    p.text = "Feed 2: The AWS Bill"
    p.level = 1
    p = tf4.add_paragraph()
    p.text = "Actual $ Chandra costs to run (Feeds Cost / C)."
    p.level = 2

    p = tf4.add_paragraph()
    p.text = "Feed 3: The Answer Key"
    p.level = 1
    p = tf4.add_paragraph()
    p.text = "Were recommendations actually correct vs ground truth? (Feeds Validation / V)."
    p.level = 2

    p = tf4.add_paragraph()
    p.text = "DPI Engine & Scoreboard:"
    p.level = 1
    p = tf4.add_paragraph()
    p.text = "Connectors pull feeds, formula computes: PI = (P × Q × 1.5E) + (G × 1.5R) + (C × V)."
    p.level = 2

    # Save to the artifacts directory
    output_path = r"C:\Users\User\.gemini\antigravity\brain\d42c9526-0eea-40b3-97d6-ee5b2265af2c\Chandra_DPI_Workflow.pptx"
    prs.save(output_path)
    print(f"Presentation saved successfully to {output_path}")

if __name__ == "__main__":
    create_presentation()
