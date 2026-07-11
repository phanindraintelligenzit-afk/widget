import pptx
from pptx import Presentation

prs = Presentation('dpi-ls-validation.pptx')

for i, slide in enumerate(prs.slides):
    print(f"--- Slide {i+1} ---")
    for j, shape in enumerate(slide.shapes):
        if hasattr(shape, "text"):
            print(f"Shape {j}: {repr(shape.text)}")
