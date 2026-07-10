import pptx
from pptx import Presentation

prs = Presentation('dpi-ls-quality.pptx')

for slide_idx in [1, 4]: # Slide 2 and Slide 5 (0-indexed)
    slide = prs.slides[slide_idx]
    print(f"--- Slide {slide_idx+1} ---")
    for i, shape in enumerate(slide.shapes):
        if shape.shape_type == 13: # 13 is msoPICTURE
            print(f"Picture {i}: Left: {shape.left}, Top: {shape.top}, Width: {shape.width}, Height: {shape.height}")
