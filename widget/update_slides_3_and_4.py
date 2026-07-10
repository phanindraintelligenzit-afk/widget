import pptx
from pptx import Presentation
import os

prs = Presentation('dpi-ls-quality.pptx')

# --- Slide 3 ---
# Shape 2 (Text): Replace all text
slide3 = prs.slides[2]
shape2 = slide3.shapes[2]
if hasattr(shape2, "text_frame"):
    new_text_3_2 = """AgentOps
Runtime Observability
Session Replay
Token Usage
Cost
Execution Timeline
|
LangSmith
Execution Traces
Prompt Flow
LangGraph
Datasets
Experiments
|
Ragas
Answer Relevancy
Faithfulness
Context Precision
Context Recall
Correctness"""
    # Delete all paragraphs except the first one
    for i in range(len(shape2.text_frame.paragraphs)-1, 0, -1):
        p = shape2.text_frame.paragraphs[i]._element
        p.getparent().remove(p)
    # Clear runs in the first paragraph and set text
    if len(shape2.text_frame.paragraphs) > 0:
        p0 = shape2.text_frame.paragraphs[0]
        # remove all runs except first
        for i in range(len(p0.runs)-1, 0, -1):
            r = p0.runs[i]._element
            r.getparent().remove(r)
        if len(p0.runs) > 0:
            p0.runs[0].text = new_text_3_2

# Shape 3 (Table)
shape3 = slide3.shapes[3]
if shape3.has_table:
    # Row 1
    shape3.table.rows[1].cells[1].text = "LLM Runtime Observability"
    shape3.table.rows[1].cells[2].text = "• Prompt Tokens\n• Completion Tokens\n• Total Tokens\n• LLM Calls\n• Session Replay\n• Agent Actions\n• Runtime Duration\n• Cost Tracking\n• Execution Timeline"
    shape3.table.rows[1].cells[3].text = "AgentOps SDK"
    shape3.table.rows[1].cells[4].text = "https://app.agentops.ai"

    # Row 2
    shape3.table.rows[2].cells[1].text = "LLM Execution Tracing"
    shape3.table.rows[2].cells[2].text = "• Prompt Flow\n• Chain Execution\n• LangGraph Nodes\n• Dataset Runs\n• Trace Hierarchy\n• Span Execution\n• Prompt Versions\n• Experiment Tracking"
    shape3.table.rows[2].cells[3].text = "LangSmith"
    shape3.table.rows[2].cells[4].text = "https://smith.langchain.com"

    # Row 3
    shape3.table.rows[3].cells[1].text = "LLM Quality Evaluation"
    shape3.table.rows[3].cells[2].text = "• Answer Relevancy\n• Faithfulness\n• Context Precision\n• Context Recall\n• Answer Correctness\n• Semantic Similarity\n• Response Groundedness"
    shape3.table.rows[3].cells[3].text = "Ragas"
    shape3.table.rows[3].cells[4].text = "https://docs.ragas.io"


# --- Slide 4 ---
slide4 = prs.slides[3]

# Shape 2 (Picture)
pic_shape = slide4.shapes[2]
if pic_shape.shape_type == 13: # Picture
    left = pic_shape.left
    top = pic_shape.top
    width = pic_shape.width
    height = pic_shape.height
    
    # remove old picture
    sp = pic_shape.element
    sp.getparent().remove(sp)
    
    # add new picture
    slide4.shapes.add_picture('quality_example.png', left, top, width, height)

# Shape 3 (Text block on right side)
shape_4_3 = slide4.shapes[3]
if hasattr(shape_4_3, "text_frame"):
    new_text_4_3 = """Formula
Quality Score
=
Successful Quality Checks
÷
Total Required Quality Checks

Explain
If all quality metrics pass
Score = 1.0
If several metrics fail
Score decreases.

Good Example
Passed
10/10
Score
1.0

Bad Example
Passed
5/10
Score
0.5

Safety Rule
If the Quality Score falls below the configured threshold
The overall DPI-LS score is reduced
The execution is flagged
Manual review is required before production deployment."""
    # Delete all paragraphs except the first one
    for i in range(len(shape_4_3.text_frame.paragraphs)-1, 0, -1):
        p = shape_4_3.text_frame.paragraphs[i]._element
        p.getparent().remove(p)
    # Clear runs in the first paragraph and set text
    if len(shape_4_3.text_frame.paragraphs) > 0:
        p0 = shape_4_3.text_frame.paragraphs[0]
        # remove all runs except first
        for i in range(len(p0.runs)-1, 0, -1):
            r = p0.runs[i]._element
            r.getparent().remove(r)
        if len(p0.runs) > 0:
            p0.runs[0].text = new_text_4_3


prs.save('dpi-ls-quality-temp.pptx')
print("Successfully updated Slides 3 and 4 in dpi-ls-quality.pptx")
