import pptx
from pptx import Presentation
import os

prs = Presentation('dpi-ls-quality.pptx')

media_dir = r"C:\Users\User\.gemini\antigravity\brain\2b74a66f-02ca-47cf-8cf2-341f2c230b3f"
# The exact filenames are dynamic, let's list them
files = os.listdir(media_dir)
pngs = [f for f in files if f.endswith('.png') and 'media__' in f]
pngs.sort(reverse=True) # newest first

if len(pngs) >= 2:
    img1 = os.path.join(media_dir, pngs[0])
    img2 = os.path.join(media_dir, pngs[1])
else:
    img1 = None
    img2 = None

def replace_picture(slide, pic_shape, new_img_path):
    left = pic_shape.left
    top = pic_shape.top
    width = pic_shape.width
    height = pic_shape.height
    
    # remove old picture
    sp = pic_shape.element
    sp.getparent().remove(sp)
    
    # add new picture
    slide.shapes.add_picture(new_img_path, left, top, width, height)

if img1 and img2:
    # Slide 2
    slide2 = prs.slides[1]
    pics2 = [s for s in slide2.shapes if s.shape_type == 13]
    if pics2:
        replace_picture(slide2, pics2[0], img1)
        
    # Slide 5
    slide5 = prs.slides[4]
    pics5 = [s for s in slide5.shapes if s.shape_type == 13]
    if len(pics5) >= 1:
        replace_picture(slide5, pics5[0], img1) # Replace Left Image
    if len(pics5) >= 2:
        replace_picture(slide5, pics5[1], img2) # Replace Right Top Image
    if len(pics5) >= 3:
        replace_picture(slide5, pics5[2], img1) # Replace Right Middle Image
    if len(pics5) >= 4:
        # Just remove the 4th picture if we don't have enough distinct ones,
        # or reuse img2
        replace_picture(slide5, pics5[3], img2)

prs.save('dpi-ls-quality.pptx')
print("Successfully replaced images in dpi-ls-quality.pptx")
