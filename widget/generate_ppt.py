import pptx
from pptx import Presentation
import sys
import shutil
import os

shutil.copy('dpi-ls-validation.pptx', 'dpi-ls-quality.pptx')

prs = Presentation('dpi-ls-quality.pptx')

# We'll do a simple find/replace on all text runs
replacements = {
    # Slide 1
    "VALIDATION DIMENSION": "QUALITY DIMENSION",
    "Cost Traceability, Real-Time Telemetry & Core Scoring Validation": "Answer Quality, Trustworthiness, LLM Evaluation & Runtime Quality Scoring",
    # Slide 2
    "What is DPI-LS Validation & How It Works": "What is DPI-LS Quality & How It Works",
    "Digital FTE Validation Index": "Digital FTE Quality Index",
    "Validation (V) measures whether an AI agent produces trustworthy, accurate, observable and verifiable outputs during runtime.": "Quality (Q) measures whether an AI Agent produces accurate, relevant, trustworthy and consistent answers during runtime.",
    "The Validation dimension combines three runtime validation resources:": "Quality combines three runtime evaluation resources:",
    "DeepEval": "AgentOps",
    "Jaeger": "LangSmith",
    "Zipkin": "Ragas",
    "AI answer quality": "LLM answer quality",
    "Runtime correctness": "Prompt execution",
    "Distributed execution tracing": "Agent workflow\n   Runtime observability",
    "End-to-end observability": "Retrieval quality\n   AI reasoning consistency",
    "Validation ensures every AI execution can be evaluated, traced and verified before production deployment.": "Quality ensures every AI execution can be evaluated before production deployment.",
    "Runtime telemetry is collected from DeepEval, Jaeger, and Zipkin using real OpenTelemetry traces in the development environment.": "Runtime telemetry is collected from AgentOps, LangSmith, and Ragas using real runtime execution.",
    
    # Slide 3
    "Validation (V) Dimension & Telemetry Sources": "Quality (Q) Dimension & Telemetry Sources",
    "DeepEval: Answer Relevancy, Faithfulness, Hallucination, Correctness | Jaeger: Trace ID, Span Count, Latency, Dependencies | Zipkin: Trace Timeline, Execution Timeline, Service Calls, Trace Latency": 
    "1 | LLM Runtime Observability | Captures: LLM execution, Token usage, Sessions, Actions, Cost, Runtime telemetry | AgentOps SDK | https://app.agentops.ai\n2 | LLM Execution Tracing | Captures: Prompt execution, Chains, Runs, Datasets, Experiments, Execution traces | LangSmith | https://smith.langchain.com\n3 | LLM Evaluation Metrics | Evaluates: Answer Relevancy, Faithfulness, Context Precision, Context Recall, Answer Correctness | Ragas | https://docs.ragas.io",
    
    # Slide 4
    "Formula, Examples & Validation Rule": "Formula, Examples & Quality Rule",
    "Validation Score = Correct Validation Checks \u00f7 Total Validation Checks": "Quality Score = Successful Quality Checks \u00f7 Total Required Quality Checks",
    "Validation Score = Correct Validation Checks \u00f7 Total Validation Checks": "Quality Score = Successful Quality Checks \u00f7 Total Required Quality Checks",
    "Validation Score": "Quality Score",
    "Validation Checks": "Quality Checks",
    "Validation": "Quality",
    "Quality score": "Quality Score",
    "Quality Rule": "Quality Rule",
    
    # Slide 5
    "DPI-LS Validation Dashboards & Runtime Evidence": "DPI-LS Quality Dashboards & Runtime Evidence",
    "Validation Resource Evaluation Dashboard": "Quality Resources Dashboard",
    "Validation dashboards showing real runtime evaluation and distributed observability.": "Quality dashboards showing live runtime quality evaluation and observability.",
    "Evidence shown is from live runtime telemetry captured in the development environment using OpenTelemetry.": "Evidence is collected from AgentOps, LangSmith and Ragas during development using OpenTelemetry.",
    
    # Slide 6
    "Production Validation Execution Workflow": "Production Quality Execution Workflow",
    "Start Backend": "Start Backend\nuv run uvicorn api.app:app --host 127.0.0.1 --port 8000 --reload",
    "Run Validation Agent": "Run Quality Agent\nuv run python examples/test_agent.py",
    "Execute Validation Evaluation": "Execute Quality Evaluation",
    "Verify Results": "Verify Results",
    "Live Runtime Metrics Displayed": "Live Runtime Metrics Displayed",
    
    # Slide 7
    "Validation Dimension Summary": "Quality Dimension Summary",
    "DeepEval Integrated": "AgentOps Integrated",
    "Jaeger Integrated": "LangSmith Integrated",
    "Zipkin Integrated": "Ragas Integrated",
    "Validation Resources Dashboard": "Quality Resources Dashboard"
}

# Apply text replacement
for slide in prs.slides:
    for shape in slide.shapes:
        if hasattr(shape, "text_frame"):
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    original_text = run.text
                    for old_text, new_text in replacements.items():
                        if old_text in run.text:
                            run.text = run.text.replace(old_text, new_text)

# We also need to do a full-text replacement for things that might span runs
for slide in prs.slides:
    for shape in slide.shapes:
        if hasattr(shape, "text_frame"):
            for paragraph in shape.text_frame.paragraphs:
                p_text = "".join(r.text for r in paragraph.runs)
                modified = False
                for old_text, new_text in replacements.items():
                    if old_text in p_text and old_text not in [r.text for r in paragraph.runs]:
                        p_text = p_text.replace(old_text, new_text)
                        modified = True
                if modified and len(paragraph.runs) > 0:
                    for i in range(1, len(paragraph.runs)):
                        paragraph.runs[i].text = ""
                    paragraph.runs[0].text = p_text

# Additional specific text replacements
def force_replace(old, new):
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text_frame"):
                for paragraph in shape.text_frame.paragraphs:
                    p_text = "".join(r.text for r in paragraph.runs)
                    if old in p_text:
                        for i in range(1, len(paragraph.runs)):
                            paragraph.runs[i].text = ""
                        paragraph.runs[0].text = p_text.replace(old, new)

force_replace("Validation Dimension Summary", "Quality Dimension Summary")
force_replace("Production Validation Pending", "Production Validation Pending")
force_replace("Quality (V)", "Quality (Q)")
force_replace("Validation (V)", "Quality (Q)")
force_replace("Validation Rule", "Quality Rule")

# Save the updated pptx
prs.save('dpi-ls-quality.pptx')
print("Successfully generated dpi-ls-quality.pptx")
