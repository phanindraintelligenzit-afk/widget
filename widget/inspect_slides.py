import pptx
from pptx import Presentation

prs = Presentation('dpi-ls-quality.pptx')

print("--- Slide 3 ---")
for i, shape in enumerate(prs.slides[2].shapes):
    if hasattr(shape, "text"):
        print(f"Shape {i} (Text): {repr(shape.text)}")
    elif shape.shape_type == 13: # Picture
        print(f"Shape {i} (Picture)")
    elif shape.shape_type == 19: # Table
        print(f"Shape {i} (Table)")
        for row in shape.table.rows:
            row_data = [cell.text for cell in row.cells]
            print(f"  Row: {row_data}")
    elif shape.shape_type == 6: # Group
        print(f"Shape {i} (Group)")

print("\n--- Slide 4 ---")
for i, shape in enumerate(prs.slides[3].shapes):
    if hasattr(shape, "text"):
        print(f"Shape {i} (Text): {repr(shape.text)}")
    elif shape.shape_type == 13: # Picture
        print(f"Shape {i} (Picture)")
    elif shape.shape_type == 19: # Table
        print(f"Shape {i} (Table)")
    elif shape.shape_type == 6: # Group
        print(f"Shape {i} (Group)")

