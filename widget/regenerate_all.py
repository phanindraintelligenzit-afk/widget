import pptx
from pptx import Presentation
from pptx.util import Pt, Inches
import shutil
import os

# Clean up previous
for f in ['dpi-ls-quality.pptx', 'dpi-ls-quality-temp.pptx', 'dpi-ls-quality-final.pptx']:
    if os.path.exists(f):
        try:
            os.remove(f)
        except:
            pass

shutil.copy('dpi-ls-validation.pptx', 'dpi-ls-quality.pptx')

prs = Presentation('dpi-ls-quality.pptx')

replacements = {
    # Slide 1
    "VALIDATION DIMENSION": "QUALITY DIMENSION",
    "Cost Traceability, Real-Time Telemetry & Core Scoring Validation": "Answer Quality, Trustworthiness, LLM Evaluation & Runtime Quality Scoring",
    # Slide 2
    "What is DPI-LS Validation & How It Works": "What is DPI-LS Quality & How It Works",
    "Digital FTE Validation Index": "Digital FTE Quality Index",
    "Validation (V) measures whether an AI agent produces trustworthy, accurate, observable and verifiable outputs during runtime.": "Quality (Q) measures whether an AI Agent produces accurate, relevant, trustworthy and consistent answers during runtime.",
    "The Validation dimension combines three runtime validation resources:": "Quality combines three runtime evaluation resources:",
    "DeepEval": "AgentOps",
    "Jaeger": "LangSmith",
    "Zipkin": "Ragas",
    "AI answer quality": "LLM answer quality",
    "Runtime correctness": "Prompt execution",
    "Distributed execution tracing": "Agent workflow\n   Runtime observability",
    "End-to-end observability": "Retrieval quality\n   AI reasoning consistency",
    "Validation ensures every AI execution can be evaluated, traced and verified before production deployment.": "Quality ensures every AI execution can be evaluated before production deployment.",
    "Runtime telemetry is collected from DeepEval, Jaeger, and Zipkin using real OpenTelemetry traces in the development environment.": "Runtime telemetry is collected from AgentOps, LangSmith, and Ragas using real runtime execution.",
    
    # Slide 3 generic replacements
    "Validation (V) Dimension & Telemetry Sources": "Quality (Q) Dimension & Telemetry Sources",
    
    # Slide 4 generic replacements
    "Formula, Examples & Validation Rule": "Formula, Examples & Quality Rule",
    "Validation Score": "Quality Score",
    "Validation Checks": "Quality Checks",
    "Validation": "Quality",
    "Quality Rule": "Quality Rule",
    
    # Slide 5
    "DPI-LS Validation Dashboards & Runtime Evidence": "DPI-LS Quality Dashboards & Runtime Evidence",
    "Validation Resource Evaluation Dashboard": "Quality Resources Dashboard",
    "Validation dashboards showing real runtime evaluation and distributed observability.": "Quality dashboards showing live runtime quality evaluation and observability.",
    "Evidence shown is from live runtime telemetry captured in the development environment using OpenTelemetry.": "Evidence is collected from AgentOps, LangSmith and Ragas during development using OpenTelemetry.",
    
    # Slide 6
    "Production Validation Execution Workflow": "Production Quality Execution Workflow",
    "Start Backend": "Start Backend\nuv run uvicorn api.app:app --host 127.0.0.1 --port 8000 --reload",
    "Run Validation Agent": "Run Quality Agent\nuv run python examples/test_agent.py",
    "Execute Validation Evaluation": "Execute Quality Evaluation",
    "Verify Results": "Verify Results",
    "Live Runtime Metrics Displayed": "Live Runtime Metrics Displayed",
    
    # Slide 7
    "Validation Dimension Summary": "Quality Dimension Summary",
    "Validation Resources Dashboard": "Quality Resources Dashboard"
}

# Apply text replacement
for slide in prs.slides:
    for shape in slide.shapes:
        if hasattr(shape, "text_frame"):
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    for old_text, new_text in replacements.items():
                        if old_text in run.text:
                            run.text = run.text.replace(old_text, new_text)

# Full-text replacement for things that might span runs
for slide in prs.slides:
    for shape in slide.shapes:
        if hasattr(shape, "text_frame"):
            for paragraph in shape.text_frame.paragraphs:
                p_text = "".join(r.text for r in paragraph.runs)
                modified = False
                for old_text, new_text in replacements.items():
                    if old_text in p_text and old_text not in [r.text for r in paragraph.runs]:
                        p_text = p_text.replace(old_text, new_text)
                        modified = True
                if modified and len(paragraph.runs) > 0:
                    for i in range(1, len(paragraph.runs)):
                        paragraph.runs[i].text = ""
                    paragraph.runs[0].text = p_text

# Additional specific text replacements
def force_replace(old, new):
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text_frame"):
                for paragraph in shape.text_frame.paragraphs:
                    p_text = "".join(r.text for r in paragraph.runs)
                    if old in p_text:
                        for i in range(1, len(paragraph.runs)):
                            paragraph.runs[i].text = ""
                        paragraph.runs[0].text = p_text.replace(old, new)

force_replace("Validation Dimension Summary", "Quality Dimension Summary")
force_replace("Production Validation Pending", "Production Validation Pending")
force_replace("Quality (V)", "Quality (Q)")
force_replace("Validation (V)", "Quality (Q)")
force_replace("Validation Rule", "Quality Rule")
force_replace("DeepEval Integrated", "AgentOps Integrated")
force_replace("Jaeger Integrated", "LangSmith Integrated")
force_replace("Zipkin Integrated", "Ragas Integrated")

# --- Slide 3 specific fixes ---
slide3 = prs.slides[2]

# Fix 1: Shape 2 (Runtime Metrics description) to ONLY ONE LINE
shape2 = slide3.shapes[2]
if hasattr(shape2, "text_frame"):
    new_metrics_text = "AgentOps \u2013 Runtime Observability, Tokens, Sessions, Cost | LangSmith \u2013 Traces, Prompt Flow, Chains, Experiments | Ragas \u2013 Relevancy, Faithfulness, Precision, Recall, Correctness"
    # Keep the original formatting of the first run
    if len(shape2.text_frame.paragraphs) > 0:
        p0 = shape2.text_frame.paragraphs[0]
        for i in range(len(shape2.text_frame.paragraphs)-1, 0, -1):
            p = shape2.text_frame.paragraphs[i]._element
            p.getparent().remove(p)
        for i in range(len(p0.runs)-1, 0, -1):
            r = p0.runs[i]._element
            r.getparent().remove(r)
        if len(p0.runs) > 0:
            p0.runs[0].text = new_metrics_text
            p0.runs[0].font.size = Pt(16) # Explicitly set a reasonable size so it fits on one line
        else:
            p0.text = new_metrics_text
            p0.font.size = Pt(16)

# Fix 2: Table contents, margins, font sizes
shape3 = slide3.shapes[3]
if shape3.has_table:
    # Set slightly larger row heights and adjust position if needed
    # The user says "Increase table height slightly", let's adjust row heights.
    # We won't change the overall width or x,y position drastically to keep it looking the same.
    for row in shape3.table.rows:
        row.height = int(row.height * 1.2)
        
    table_data = [
        # Row 1
        ("LLM Runtime Observability", "\u2022 Tokens\n\u2022 LLM Calls\n\u2022 Sessions\n\u2022 Runtime\n\u2022 Cost\n\u2022 Timeline", "AgentOps SDK", "https://app.agentops.ai"),
        # Row 2
        ("LLM Execution Tracing", "\u2022 Prompt Flow\n\u2022 Chains\n\u2022 LangGraph\n\u2022 Datasets\n\u2022 Experiments\n\u2022 Traces", "LangSmith", "https://smith.langchain.com"),
        # Row 3
        ("LLM Quality Evaluation", "\u2022 Relevancy\n\u2022 Faithfulness\n\u2022 Context Precision\n\u2022 Context Recall\n\u2022 Correctness", "Ragas", "https://docs.ragas.io")
    ]
    
    for row_idx, data in enumerate(table_data, start=1):
        cells = shape3.table.rows[row_idx].cells
        
        # Parameter
        cells[1].text = data[0]
        # What It Measures
        cells[2].text = data[1]
        # Source System
        cells[3].text = data[2]
        # URL
        cells[4].text = data[3]
        
        for col_idx in range(5):
            cell = cells[col_idx]
            # Reduce cell padding
            cell.margin_top = Inches(0.02)
            cell.margin_bottom = Inches(0.02)
            cell.margin_left = Inches(0.05)
            cell.margin_right = Inches(0.05)
            
            # Reduce font size to approximately 15-16 pt, enable word wrapping
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(15.5)
            cell.text_frame.word_wrap = True

# --- Slide 4 specific fixes ---
slide4 = prs.slides[3]

# Shape 2 (Picture) - Replace with new Quality Example image
pic_shape = slide4.shapes[2]
if pic_shape.shape_type == 13: # Picture
    left = pic_shape.left
    top = pic_shape.top
    width = pic_shape.width
    height = pic_shape.height
    sp = pic_shape.element
    sp.getparent().remove(sp)
    if os.path.exists('quality_example.png'):
        slide4.shapes.add_picture('quality_example.png', left, top, width, height)

# Shape 3 (Text block on right side)
shape_4_3 = slide4.shapes[3]
if hasattr(shape_4_3, "text_frame"):
    new_text_4_3 = """Formula
Quality Score
=
Successful Quality Checks
\u00f7
Total Required Quality Checks

Explain
If all quality metrics pass
Score = 1.0
If several metrics fail
Score decreases.

Good Example
Passed
10/10
Score
1.0

Bad Example
Passed
5/10
Score
0.5

Safety Rule
If the Quality Score falls below the configured threshold
The overall DPI-LS score is reduced
The execution is flagged
Manual review is required before production deployment."""
    for i in range(len(shape_4_3.text_frame.paragraphs)-1, 0, -1):
        p = shape_4_3.text_frame.paragraphs[i]._element
        p.getparent().remove(p)
    if len(shape_4_3.text_frame.paragraphs) > 0:
        p0 = shape_4_3.text_frame.paragraphs[0]
        for i in range(len(p0.runs)-1, 0, -1):
            r = p0.runs[i]._element
            r.getparent().remove(r)
        if len(p0.runs) > 0:
            p0.runs[0].text = new_text_4_3

prs.save('dpi-ls-quality.pptx')
print("Successfully regenerated dpi-ls-quality.pptx")
