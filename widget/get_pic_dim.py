import pptx
from pptx import Presentation

prs = Presentation('dpi-ls-quality.pptx')
pic = prs.slides[3].shapes[2]
print(f"Slide 4 Picture: left={pic.left}, top={pic.top}, width={pic.width}, height={pic.height}")
